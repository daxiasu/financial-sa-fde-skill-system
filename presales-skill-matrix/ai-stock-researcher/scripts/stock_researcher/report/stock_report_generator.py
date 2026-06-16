#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
股票综合研究报告生成器
Stock Research Report Generator

功能：
1. 综合市场全量股票分析
2. 东方财富/雪球股吧情绪分析
3. 历史案例分析
4. 量化技术分析
5. 短中期预测
6. PPT/PDF导出

导出格式：
- PPT: 可视化演示文稿
- PDF: 详细研究报告
"""

import json
import os
import time
import ssl
import re
import urllib.request
from typing import List, Dict, Optional, Tuple
from pathlib import Path
from datetime import datetime, timedelta
from collections import Counter

# PPT生成
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# PDF生成
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image
    from reportlab.lib import colors
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# 尝试导入crawl_utils
try:
    from crawl_utils import safe_request, detect_encoding, fetch_json
    HAS_CRAWL_UTILS = True
except ImportError:
    HAS_CRAWL_UTILS = False


class StockResearchReportGenerator:
    """
    股票综合研究报告生成器

    生成步骤：
    1. 股票基本信息（代码、名称、行业、市值）
    2. 实时行情与技术分析
    3. 东方财富/雪球情绪分析
    4. 历史案例匹配
    5. 量化多维分析
    6. 短中期预测
    7. 投资建议
    8. 导出PPT/PDF
    """

    def __init__(self, data_dir: str = None, output_dir: str = None):
        if data_dir is None:
            data_dir = Path(__file__).resolve().parents[2] / "data"
        if output_dir is None:
            output_dir = Path(__file__).resolve().parents[2] / "data" / "reports"

        self.data_dir = Path(data_dir)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        self.ctx = ssl.create_default_context()
        self.ctx.check_hostname = False
        self.ctx.verify_mode = ssl.CERT_NONE

    def fetch_stock_data(self, code: str) -> Dict:
        """
        获取股票综合数据

        Args:
            code: 股票代码

        Returns:
            Dict: 综合股票数据
        """
        code = str(code).zfill(6)

        # 实时行情
        quote = self._fetch_realtime_quote(code)

        # K线数据
        kline = self._fetch_kline(code, days=120)

        # 基本信息
        info = self._fetch_stock_info(code)

        # 技术指标
        tech = self._calc_technical_indicators(kline)

        return {
            "code": code,
            "name": info.get("name", code),
            "sector": info.get("sector", ""),
            "market_cap": quote.get("market_cap", 0),
            "float_capital": quote.get("float_capital", 0),
            "price": quote.get("price", 0),
            "prev_close": quote.get("prev_close", 0),
            "change_pct": quote.get("change_pct", 0),
            "volume": quote.get("volume", 0),
            "amount": quote.get("amount", 0),
            "turnover_rate": quote.get("turnover_rate", 0),
            "high": quote.get("high", 0),
            "low": quote.get("low", 0),
            "open": quote.get("open", 0),
            "kline": kline,
            "technical": tech
        }

    def _fetch_realtime_quote(self, code: str) -> Dict:
        """获取实时行情"""
        if code.startswith(("6", "5", "9")):
            ts = f"sh{code}"
        else:
            ts = f"sz{code}"

        url = f"https://qt.gtimg.cn/q={ts}&_={int(time.time()*1000)}"

        headers = {
            "User-Agent": "Mozilla/5.0",
            "Referer": "https://gu.qq.com/"
        }

        try:
            if HAS_CRAWL_UTILS:
                raw = safe_request(url, headers=headers, timeout=8)
            else:
                req = urllib.request.Request(url, headers=headers)
                with urllib.request.urlopen(req, timeout=8, context=self.ctx) as resp:
                    raw = resp.read()

            text = raw.decode("gbk", errors="replace")

            m = re.search(r'v_\w+="(.+?)"', text)
            if not m:
                return {}

            fields = m.group(1).split("~")
            if len(fields) < 45:
                return {}

            return {
                "price": float(fields[3]) if fields[3] else 0,
                "prev_close": float(fields[4]) if fields[4] else 0,
                "open": float(fields[5]) if fields[5] else 0,
                "volume": float(fields[6]) if fields[6] else 0,
                "turnover_rate": float(fields[8]) if fields[8] else 0,
                "market_cap": float(fields[44]) if fields[44] else 0,
                "float_capital": float(fields[45]) if fields[45] else 0,
                "change_pct": float(fields[31]) if fields[31] else 0,
                "high": float(fields[33]) if fields[33] else 0,
                "low": float(fields[34]) if fields[34] else 0,
                "amount": float(fields[37]) if fields[37] else 0,
            }

        except Exception as e:
            print(f"获取实时行情失败: {e}")
            return {}

    def _fetch_kline(self, code: str, days: int = 120) -> Dict:
        """获取K线数据"""
        if code.startswith(("6", "5", "9")):
            prefix = "sh"
        else:
            prefix = "sz"

        url = f"https://web.ifzq.gtimg.cn/appstock/app/fqkline/get?_var=kline_dayqfq&param={prefix}{code},day,,,{days},qfq&r=0.1"

        headers = {
            "User-Agent": "Mozilla/5.0",
            "Referer": "https://gu.qq.com/"
        }

        try:
            if HAS_CRAWL_UTILS:
                raw = safe_request(url, headers=headers, timeout=8)
                if isinstance(raw, tuple):
                    raw = raw[0]
            else:
                req = urllib.request.Request(url, headers=headers)
                with urllib.request.urlopen(req, timeout=8, context=self.ctx) as resp:
                    raw = resp.read()

            text = raw.decode("utf-8", errors="replace")

            matches = re.findall(
                r'\["(\d{4}-\d{2}-\d{2})",\s*"([\d.]+)",\s*"([\d.]+)",\s*"([\d.]+)",\s*"([\d.]+)",\s*"([\d.]+)"\]',
                text
            )

            if not matches:
                return {}

            dates, opens, highs, lows, closes, volumes = [], [], [], [], [], []
            for m in matches:
                dates.append(m[0])
                opens.append(float(m[1]))
                highs.append(float(m[2]))
                lows.append(float(m[3]))
                closes.append(float(m[4]))
                volumes.append(float(m[5]))

            return {
                "dates": dates,
                "opens": opens,
                "highs": highs,
                "lows": lows,
                "closes": closes,
                "volumes": volumes
            }

        except Exception as e:
            print(f"获取K线失败: {e}")
            return {}

    def _fetch_stock_info(self, code: str) -> Dict:
        """获取股票基本信息"""
        return {"name": "", "sector": ""}

    def _calc_technical_indicators(self, kline: Dict) -> Dict:
        """计算技术指标"""
        if not kline or not kline.get("closes"):
            return {}

        closes = kline["closes"]
        volumes = kline.get("volumes", [])
        highs = kline.get("highs", [])
        lows = kline.get("lows", [])

        n = len(closes)
        if n < 20:
            return {}

        # 移动平均线
        ma5 = sum(closes[-5:]) / 5 if n >= 5 else sum(closes) / n
        ma10 = sum(closes[-10:]) / 10 if n >= 10 else ma5
        ma20 = sum(closes[-20:]) / 20 if n >= 20 else ma10
        ma60 = sum(closes[-60:]) / 60 if n >= 60 else ma20

        current_price = closes[-1]

        # 均线排列
        if ma5 > ma10 > ma20 > ma60:
            ma_arrangement = "多头排列"
        elif ma5 < ma10 < ma20 < ma60:
            ma_arrangement = "空头排列"
        else:
            ma_arrangement = "震荡整理"

        # RSI
        def calc_rsi(prices, period=14):
            if len(prices) < period + 1:
                return 50
            gains = []
            losses = []
            for i in range(1, min(len(prices), period + 1)):
                diff = prices[i] - prices[i-1]
                if diff > 0:
                    gains.append(diff)
                else:
                    losses.append(abs(diff))
            avg_gain = sum(gains) / period if gains else 0
            avg_loss = sum(losses) / period if losses else 0
            if avg_loss == 0:
                return 100
            rs = avg_gain / avg_loss
            return 100 - (100 / (1 + rs))

        rsi6 = calc_rsi(closes, 6)
        rsi14 = calc_rsi(closes, 14)

        # MACD
        def calc_ema(prices, period):
            if not prices:
                return 0
            ema = prices[0]
            multiplier = 2 / (period + 1)
            for price in prices[1:]:
                ema = (price - ema) * multiplier + ema
            return ema

        ema12 = calc_ema(closes, 12)
        ema26 = calc_ema(closes, 26)
        dif = ema12 - ema26
        dea = calc_ema([dif], 9)
        macd_hist = (dif - dea) * 2

        # 布林带
        def calc_bollinger(prices, period=20, std_dev=2):
            if len(prices) < period:
                return 0, 0, 0
            recent = prices[-period:]
            sma = sum(recent) / period
            variance = sum((p - sma) ** 2 for p in recent) / period
            std = variance ** 0.5
            return sma + std_dev * std, sma, sma - std_dev * std

        upper, middle, lower = calc_bollinger(closes)

        # 技术信号
        signal = "中性"
        if ma_arrangement == "多头排列" and rsi14 < 70:
            signal = "谨慎看多"
        elif ma_arrangement == "空头排列" and rsi14 > 30:
            signal = "谨慎看空"
        elif rsi14 > 70:
            signal = "超买"
        elif rsi14 < 30:
            signal = "超卖"

        return {
            "ma5": round(ma5, 2),
            "ma10": round(ma10, 2),
            "ma20": round(ma20, 2),
            "ma60": round(ma60, 2),
            "ma_arrangement": ma_arrangement,
            "rsi6": round(rsi6, 1),
            "rsi14": round(rsi14, 1),
            "dif": round(dif, 4),
            "dea": round(dea, 4),
            "macd_hist": round(macd_hist, 4),
            "bollinger_upper": round(upper, 2),
            "bollinger_middle": round(middle, 2),
            "bollinger_lower": round(lower, 2),
            "signal": signal
        }

    def analyze_sentiment(self, code: str) -> Dict:
        """情绪分析（简化版，实际应调用sentiment_forum_crawler）"""
        # 这里应该调用 sentiment_forum_crawler 的分析功能
        # 为了完整性，这里提供默认结构
        return {
            "platform": "东方财富+雪球",
            "total_posts": 0,
            "bullish_ratio": 50,
            "sentiment_score": 0,
            "sentiment_label": "中性",
            "key_topics": [],
            "alerts": []
        }

    def analyze_historical_cases(self, code: str, kline: Dict) -> Dict:
        """
        历史案例分析

        查找相似走势的股票和历史行情
        """
        if not kline or not kline.get("closes"):
            return {}

        closes = kline["closes"]
        if len(closes) < 20:
            return {}

        # 当前走势特征
        recent_change = (closes[-1] - closes[-20]) / closes[-20] * 100 if len(closes) >= 20 else 0

        # 查找历史相似走势（这里简化，实际应比对数据库）
        return {
            "recent_change": round(recent_change, 2),
            "similar_cases": [],
            "pattern_type": "震荡整理",
            "historical_success_rate": 0.5
        }

    def generate_prediction(self, stock_data: Dict, period: str = "medium") -> Dict:
        """
        生成预测

        Args:
            stock_data: 股票数据
            period: 预测周期 short/medium/long

        Returns:
            Dict: 预测结果
        """
        tech = stock_data.get("technical", {})
        price = stock_data.get("price", 0)

        # 简单预测模型
        if period == "short":
            # 1-5日预测
            signal = tech.get("signal", "中性")
            if signal == "超买":
                predicted = -2
            elif signal == "超卖":
                predicted = 3
            else:
                predicted = 0.5

            return {
                "period": "1-5日",
                "signal": signal,
                "predicted_return": predicted,
                "confidence": 0.5,
                "key_factors": ["技术面", "市场情绪"]
            }

        elif period == "medium":
            # 20-60日预测
            ma_arrangement = tech.get("ma_arrangement", "")
            change_pct = stock_data.get("change_pct", 0)

            if ma_arrangement == "多头排列":
                predicted = 8
                signal = "看多"
            elif ma_arrangement == "空头排列":
                predicted = -5
                signal = "看空"
            else:
                predicted = 2
                signal = "中性"

            return {
                "period": "20-60日",
                "signal": signal,
                "predicted_return": predicted,
                "confidence": 0.6,
                "p10": predicted - 5,
                "p50": predicted,
                "p90": predicted + 8,
                "key_factors": ["均线系统", "趋势力度", "市场情绪"]
            }

        else:
            # 长期预测
            return {
                "period": "120+日",
                "signal": "中性",
                "predicted_return": 0,
                "confidence": 0.4,
                "key_factors": ["基本面", "行业发展"]
            }

    def generate_comprehensive_report(self, code: str, period: str = "medium") -> Dict:
        """
        生成综合研究报告

        Args:
            code: 股票代码
            period: 报告周期 short/medium/long

        Returns:
            Dict: 报告数据
        """
        print(f"[研究报告] 开始生成 {code} 的研究报告...")

        # 获取数据
        stock_data = self.fetch_stock_data(code)

        # 情绪分析
        sentiment = self.analyze_sentiment(code)

        # 历史案例
        cases = self.analyze_historical_cases(code, stock_data.get("kline", {}))

        # 预测
        prediction = self.generate_prediction(stock_data, period)

        report = {
            "basic_info": {
                "code": code,
                "name": stock_data.get("name", code),
                "sector": stock_data.get("sector", ""),
                "price": stock_data.get("price", 0),
                "change_pct": stock_data.get("change_pct", 0),
                "prev_close": stock_data.get("prev_close", 0),
                "open": stock_data.get("open", 0),
                "high": stock_data.get("high", 0),
                "low": stock_data.get("low", 0),
                "volume": stock_data.get("volume", 0),
                "amount": stock_data.get("amount", 0),
                "turnover_rate": stock_data.get("turnover_rate", 0),
                "market_cap": stock_data.get("market_cap", 0) / 1e8,  # 转换为亿元
                "update_time": datetime.now().strftime("%Y-%m-%d %H:%M")
            },
            "technical": stock_data.get("technical", {}),
            "sentiment": sentiment,
            "historical_cases": cases,
            "prediction": prediction
        }

        print(f"[研究报告] 报告生成完成")

        return report

    def export_to_ppt(self, report: Dict, output_path: str = None) -> str:
        """
        导出PPT

        Args:
            report: 报告数据
            output_path: 输出路径

        Returns:
            str: 生成的PPT文件路径
        """
        if not PPTX_AVAILABLE:
            return "PPT导出功能需要安装 python-pptx: pip install python-pptx"

        if output_path is None:
            code = report["basic_info"]["code"]
            today = datetime.now().strftime("%Y%m%d")
            output_path = self.output_dir / f"research_report_{code}_{today}.pptx"

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 标题页
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = f"{report['basic_info']['name']}({report['basic_info']['code']})"
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = "股票综合研究报告"
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = f"生成日期: {report['basic_info']['update_time']}"
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER

        # 基本信息页
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "一、基本信息"
        p.font.size = Pt(32)
        p.font.bold = True

        content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
        tf = content.text_frame

        info = report["basic_info"]
        fields = [
            ("股票代码", info["code"]),
            ("股票名称", info["name"]),
            ("所属行业", info["sector"]),
            ("当前价格", f"{info['price']:.2f}元"),
            ("涨跌幅", f"{info['change_pct']:+.2f}%"),
            ("总市值", f"{info['market_cap']:.2f}亿元"),
            ("换手率", f"{info['turnover_rate']:.2f}%"),
        ]

        for label, value in fields:
            p = tf.add_paragraph()
            p.text = f"{label}: {value}"
            p.font.size = Pt(20)

        # 技术分析页
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "二、技术分析"
        p.font.size = Pt(32)
        p.font.bold = True

        content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
        tf = content.text_frame

        tech = report.get("technical", {})
        fields = [
            ("均线系统", f"MA5={tech.get('ma5', 0):.2f} MA10={tech.get('ma10', 0):.2f} MA20={tech.get('ma20', 0):.2f}"),
            ("均线排列", tech.get("ma_arrangement", "未知")),
            ("RSI(14)", f"{tech.get('rsi14', 50):.1f}"),
            ("MACD", f"DIF={tech.get('dif', 0):.4f} DEA={tech.get('dea', 0):.4f}"),
            ("技术信号", tech.get("signal", "中性")),
        ]

        for label, value in fields:
            p = tf.add_paragraph()
            p.text = f"{label}: {value}"
            p.font.size = Pt(18)

        # 情绪分析页
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "三、情绪分析"
        p.font.size = Pt(32)
        p.font.bold = True

        content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
        tf = content.text_frame

        sent = report.get("sentiment", {})
        fields = [
            ("数据来源", sent.get("platform", "东方财富+雪球")),
            ("帖子总数", str(sent.get("total_posts", 0))),
            ("多头比例", f"{sent.get('bullish_ratio', 50):.1f}%"),
            ("情绪评分", f"{sent.get('sentiment_score', 0):.2f}"),
            ("情绪标签", sent.get("sentiment_label", "中性")),
        ]

        for label, value in fields:
            p = tf.add_paragraph()
            p.text = f"{label}: {value}"
            p.font.size = Pt(20)

        # 预测页
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "四、投资预测"
        p.font.size = Pt(32)
        p.font.bold = True

        content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
        tf = content.text_frame

        pred = report.get("prediction", {})
        fields = [
            ("预测周期", pred.get("period", "")),
            ("预测信号", pred.get("signal", "中性")),
            ("预测涨跌幅", f"{pred.get('predicted_return', 0):+.2f}%"),
            ("置信度", f"{pred.get('confidence', 0):.0%}"),
            ("关键因素", ", ".join(pred.get("key_factors", []))),
        ]

        for label, value in fields:
            p = tf.add_paragraph()
            p.text = f"{label}: {value}"
            p.font.size = Pt(20)

        # 保存
        prs.save(str(output_path))
        print(f"[PPT生成] 已保存: {output_path}")
        return str(output_path)

    def export_to_pdf(self, report: Dict, output_path: str = None) -> str:
        """
        导出PDF

        Args:
            report: 报告数据
            output_path: 输出路径

        Returns:
            str: 生成的PDF文件路径
        """
        if not PDF_AVAILABLE:
            return "PDF导出功能需要安装 reportlab: pip install reportlab"

        if output_path is None:
            code = report["basic_info"]["code"]
            today = datetime.now().strftime("%Y%m%d")
            output_path = self.output_dir / f"research_report_{code}_{today}.pdf"

        doc = SimpleDocTemplate(str(output_path), pagesize=A4)
        elements = []

        styles = getSampleStyleSheet()

        # 标题
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=24,
            spaceAfter=20
        )

        elements.append(Paragraph(f"{report['basic_info']['name']}({report['basic_info']['code']}) - 股票综合研究报告", title_style))
        elements.append(Paragraph(f"生成日期: {report['basic_info']['update_time']}", styles['Normal']))
        elements.append(Spacer(1, 20))

        # 基本信息
        elements.append(Paragraph("一、基本信息", styles['Heading1']))
        info = report["basic_info"]
        info_table = [
            ["指标", "数值"],
            ["股票代码", info["code"]],
            ["股票名称", info["name"]],
            ["所属行业", info["sector"]],
            ["当前价格", f"{info['price']:.2f}元"],
            ["涨跌幅", f"{info['change_pct']:+.2f}%"],
            ["总市值", f"{info['market_cap']:.2f}亿元"],
            ["换手率", f"{info['turnover_rate']:.2f}%"],
        ]
        t = Table(info_table, colWidths=[150, 300])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(t)
        elements.append(Spacer(1, 20))

        # 技术分析
        elements.append(Paragraph("二、技术分析", styles['Heading1']))
        tech = report.get("technical", {})
        tech_table = [
            ["指标", "数值", "说明"],
            ["MA5", f"{tech.get('ma5', 0):.2f}", ""],
            ["MA10", f"{tech.get('ma10', 0):.2f}", ""],
            ["MA20", f"{tech.get('ma20', 0):.2f}", ""],
            ["均线排列", tech.get("ma_arrangement", "未知"), ""],
            ["RSI(14)", f"{tech.get('rsi14', 50):.1f}", "超买>70 超卖<30"],
            ["MACD", f"DIF={tech.get('dif', 0):.4f}", ""],
            ["技术信号", tech.get("signal", "中性"), ""],
        ]
        t = Table(tech_table, colWidths=[150, 200, 150])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
            ('BACKGROUND', (0, 1), (-1, -1), colors.white),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(t)
        elements.append(Spacer(1, 20))

        # 情绪分析
        elements.append(Paragraph("三、情绪分析", styles['Heading1']))
        sent = report.get("sentiment", {})
        sent_table = [
            ["指标", "数值"],
            ["数据来源", sent.get("platform", "东方财富+雪球")],
            ["帖子总数", str(sent.get("total_posts", 0))],
            ["多头比例", f"{sent.get('bullish_ratio', 50):.1f}%"],
            ["情绪评分", f"{sent.get('sentiment_score', 0):.2f}"],
            ["情绪标签", sent.get("sentiment_label", "中性")],
        ]
        t = Table(sent_table, colWidths=[150, 300])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(t)
        elements.append(Spacer(1, 20))

        # 预测
        elements.append(Paragraph("四、投资预测", styles['Heading1']))
        pred = report.get("prediction", {})
        pred_table = [
            ["指标", "数值"],
            ["预测周期", pred.get("period", "")],
            ["预测信号", pred.get("signal", "中性")],
            ["预测涨跌幅", f"{pred.get('predicted_return', 0):+.2f}%"],
            ["置信度", f"{pred.get('confidence', 0):.0%}"],
            ["关键因素", ", ".join(pred.get("key_factors", []))],
        ]
        t = Table(pred_table, colWidths=[150, 300])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(t)
        elements.append(Spacer(1, 30))

        # 风险提示
        elements.append(Paragraph("风险提示", styles['Heading2']))
        elements.append(Paragraph("1. 本报告仅供参考，不构成投资建议。", styles['Normal']))
        elements.append(Paragraph("2. 投资有风险，入市需谨慎。", styles['Normal']))
        elements.append(Paragraph("3. 请根据自身风险承受能力做出投资决策。", styles['Normal']))

        # 生成PDF
        doc.build(elements)
        print(f"[PDF生成] 已保存: {output_path}")
        return str(output_path)

    def export_report(self, code: str, period: str = "medium",
                     formats: List[str] = None) -> Dict[str, str]:
        """
        导出研究报告（多格式）

        Args:
            code: 股票代码
            period: 报告周期
            formats: 导出格式列表 ["ppt", "pdf", "txt"]

        Returns:
            Dict[str, str]: 各格式的输出路径
        """
        if formats is None:
            formats = ["ppt", "pdf"]

        # 生成报告
        report = self.generate_comprehensive_report(code, period)

        outputs = {}
        today = datetime.now().strftime("%Y%m%d")

        for fmt in formats:
            if fmt == "ppt":
                path = self.output_dir / f"research_report_{code}_{today}.pptx"
                result = self.export_to_ppt(report, path)
                outputs["ppt"] = result
            elif fmt == "pdf":
                path = self.output_dir / f"research_report_{code}_{today}.pdf"
                result = self.export_to_pdf(report, path)
                outputs["pdf"] = result
            elif fmt == "txt":
                path = self.output_dir / f"research_report_{code}_{today}.txt"
                text = self._format_report_text(report)
                with open(path, "w", encoding="utf-8") as f:
                    f.write(text)
                outputs["txt"] = str(path)
                print(f"[TXT生成] 已保存: {path}")

        return outputs

    def _format_report_text(self, report: Dict) -> str:
        """格式化文本报告"""
        lines = []
        lines.append("=" * 60)
        lines.append(f"  {report['basic_info']['name']}({report['basic_info']['code']}) 股票综合研究报告")
        lines.append("=" * 60)
        lines.append("")

        # 基本信息
        lines.append("一、基本信息")
        info = report["basic_info"]
        for key in ["code", "name", "sector", "price", "change_pct", "market_cap", "turnover_rate"]:
            lines.append(f"  {key}: {info.get(key, '')}")
        lines.append("")

        # 技术分析
        lines.append("二、技术分析")
        tech = report.get("technical", {})
        for key in ["ma5", "ma10", "ma20", "ma_arrangement", "rsi14", "signal"]:
            lines.append(f"  {key}: {tech.get(key, '')}")
        lines.append("")

        # 情绪分析
        lines.append("三、情绪分析")
        sent = report.get("sentiment", {})
        for key in ["platform", "total_posts", "bullish_ratio", "sentiment_label"]:
            lines.append(f"  {key}: {sent.get(key, '')}")
        lines.append("")

        # 预测
        lines.append("四、投资预测")
        pred = report.get("prediction", {})
        for key in ["period", "signal", "predicted_return", "confidence", "key_factors"]:
            val = pred.get(key, '')
            if isinstance(val, list):
                val = ", ".join(val)
            lines.append(f"  {key}: {val}")
        lines.append("")

        lines.append("=" * 60)
        lines.append("风险提示：以上内容仅供参考，不构成投资建议")
        lines.append("=" * 60)

        return "\n".join(lines)


def main():
    """测试"""
    print("=" * 60)
    print("  股票综合研究报告生成器 - 测试")
    print("=" * 60)

    generator = StockResearchReportGenerator()

    # 测试股票
    test_code = "600519"  # 贵州茅台

    print(f"\n生成 {test_code} 的研究报告...")

    # 导出多格式
    outputs = generator.export_report(test_code, period="medium", formats=["ppt", "pdf", "txt"])

    print("\n生成结果:")
    for fmt, path in outputs.items():
        print(f"  {fmt}: {path}")


if __name__ == "__main__":
    main()