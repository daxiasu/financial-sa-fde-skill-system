#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
股票/基金持仓导入引擎 v1.0
支持: 截图OCR识别 | PPT/Word/PDF导入 | 浏览器爬取 | Excel/CSV导出

集成 pkg/portfolio_tracker.py (SQLite) 和 tracker/portfolio.py (JSON)
"""
from __future__ import annotations

import json
import os
import re
import csv
import hashlib
from datetime import datetime
from pathlib import Path
from typing import Optional

SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent.parent
DATA_DIR = SKILL_DIR / "data"
CLIENTS_DIR = DATA_DIR / "clients"


class StockImportEngine:
    """股票/基金持仓导入引擎"""

    # 6位代码模式 (股票+基金共用)
    CODE_PATTERN = re.compile(r'(?<!\d)(\d{6})(?!\d)')

    # 份额/股数模式
    SHARES_PATTERN = re.compile(
        r'(?:份额|持有|买入|股数|数量)[:：\s]*(\d+(?:,\d{3})*(?:\.\d+)?)\s*(?:股|份|手|万份)?'
        r'|(\d+(?:,\d{3})*(?:\.\d+)?)\s*(?:股|份|手)'
    )

    # 成本模式
    COST_PATTERN = re.compile(
        r'(?:成本|买入价|价格|均价)[:：\s]*(\d+\.?\d*)'
    )

    # 日期模式
    DATE_PATTERN = re.compile(
        r'(\d{4})[-/年](\d{1,2})[-/月](\d{1,2})[日号]?'
    )

    # 金额模式
    AMOUNT_PATTERN = re.compile(
        r'(?:金额|市值|投入|总价)[:：\s]*(\d+(?:,\d{3})*(?:\.\d+)?)\s*(?:元|万元|万)?'
    )

    # 止盈止损模式
    STOP_LOSS_PATTERN = re.compile(r'(?:止损|亏损)[:：\s]*[-]?(\d+\.?\d*)\s*%?')
    TAKE_PROFIT_PATTERN = re.compile(r'(?:止盈|目标|盈利)[:：\s]*(\d+\.?\d*)\s*%?')

    def __init__(self, data_dir=None):
        self.data_dir = Path(data_dir) if data_dir else DATA_DIR
        self.clients_dir = self.data_dir / "clients"
        self.clients_dir.mkdir(parents=True, exist_ok=True)

        # 延迟初始化 tracker
        self._sqlite_tracker = None
        self._json_tracker = None

    @property
    def sqlite_tracker(self):
        """SQLite持仓跟踪器 (pkg/portfolio_tracker.py)"""
        if self._sqlite_tracker is None:
            import sys
            sys.path.insert(0, str(SKILL_DIR / "pkg"))
            try:
                from portfolio_tracker import PortfolioTracker
                db_path = self.data_dir / "portfolio.db"
                self._sqlite_tracker = PortfolioTracker(db_path=db_path)
            except Exception:
                self._sqlite_tracker = None
        return self._sqlite_tracker

    @property
    def json_tracker(self):
        """JSON持仓跟踪器 (tracker/portfolio.py)"""
        if self._json_tracker is None:
            try:
                from tracker.portfolio import PortfolioTracker as JsonTracker
                self._json_tracker = JsonTracker(data_dir=str(self.data_dir))
            except Exception:
                self._json_tracker = None
        return self._json_tracker

    # ==================== 截图OCR识别 ====================

    def import_from_screenshot(self, image_path: str, client_id: str = None) -> dict:
        """从截图OCR识别股票/基金持仓"""
        result = {'success': False, 'items': [], 'raw_text': '', 'errors': []}

        if not os.path.exists(image_path):
            result['errors'].append(f'文件不存在: {image_path}')
            return result

        text = self._ocr_with_tesseract(image_path)
        if not text:
            text = self._ocr_with_easyocr(image_path)

        if not text:
            result['errors'].append('OCR识别失败，请检查图片清晰度。')
            return result

        result['raw_text'] = text
        items = self._parse_items_from_text(text)
        result['items'] = items
        result['success'] = len(items) > 0

        if not items:
            result['errors'].append('未能从截图识别出持仓信息。')

        if client_id and items:
            self.save_to_repository(client_id, items, source='screenshot',
                                    source_path=image_path)

        return result

    def _ocr_with_tesseract(self, image_path: str) -> str:
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(image_path)
            text = pytesseract.image_to_string(img, lang='chi_sim+eng')
            return text.strip()
        except ImportError:
            return ""
        except Exception:
            return ""

    def _ocr_with_easyocr(self, image_path: str) -> str:
        try:
            import easyocr
            reader = easyocr.Reader(['ch_sim', 'en'], gpu=False)
            results = reader.readtext(image_path)
            return '\n'.join([r[1] for r in results])
        except ImportError:
            return ""
        except Exception:
            return ""

    # ==================== PPT导入 ====================

    def import_from_ppt(self, file_path: str, client_id: str = None) -> dict:
        """从PPT幻灯片导入持仓信息"""
        result = {'success': False, 'items': [], 'errors': []}

        if not os.path.exists(file_path):
            result['errors'].append(f'文件不存在: {file_path}')
            return result

        try:
            from pptx import Presentation
        except ImportError:
            result['errors'].append('请安装 python-pptx: pip install python-pptx')
            return result

        try:
            prs = Presentation(file_path)
            all_text = []
            table_data = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    # 提取表格
                    if shape.has_table:
                        table = shape.table
                        rows_data = []
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            rows_data.append(cells)
                        table_data.extend(rows_data)

                    # 提取文本框
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                all_text.append(text)

                    # 提取组合图形中的文本
                    if shape.has_text_frame:
                        all_text.append(shape.text_frame.text.strip())

            full_text = '\n'.join(all_text)

            items = []
            if table_data:
                items = self._parse_items_from_table(table_data)

            text_items = self._parse_items_from_text(full_text)
            items = self._merge_items(items, text_items)

            result['items'] = items
            result['success'] = len(items) > 0

            if not items:
                result['errors'].append('未能从PPT中识别出持仓信息。')

            if client_id and items:
                self.save_to_repository(client_id, items, source='ppt',
                                        source_path=file_path)

        except Exception as e:
            result['errors'].append(f'PPT解析失败: {str(e)}')

        return result

    # ==================== Word文档导入 ====================

    def import_from_docx(self, file_path: str, client_id: str = None) -> dict:
        """从Word文档导入持仓信息"""
        result = {'success': False, 'items': [], 'errors': []}

        if not os.path.exists(file_path):
            result['errors'].append(f'文件不存在: {file_path}')
            return result

        try:
            from docx import Document
        except ImportError:
            result['errors'].append('请安装 python-docx: pip install python-docx')
            return result

        try:
            doc = Document(file_path)
            all_text = []
            table_data = []

            for table in doc.tables:
                rows_data = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows_data.append(cells)
                table_data.extend(rows_data)

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    all_text.append(text)

            full_text = '\n'.join(all_text)

            items = []
            if table_data:
                items = self._parse_items_from_table(table_data)

            text_items = self._parse_items_from_text(full_text)
            items = self._merge_items(items, text_items)

            result['items'] = items
            result['success'] = len(items) > 0

            if not items:
                result['errors'].append('未能从文档中识别出持仓信息。')

            if client_id and items:
                self.save_to_repository(client_id, items, source='docx',
                                        source_path=file_path)

        except Exception as e:
            result['errors'].append(f'Word文档解析失败: {str(e)}')

        return result

    # ==================== PDF导入 ====================

    def import_from_pdf(self, file_path: str, client_id: str = None) -> dict:
        """从PDF文档导入持仓信息"""
        result = {'success': False, 'items': [], 'errors': []}

        if not os.path.exists(file_path):
            result['errors'].append(f'文件不存在: {file_path}')
            return result

        text = self._extract_pdf_text(file_path)
        if not text:
            result['errors'].append('PDF文本提取失败，文件可能为扫描件请使用截图导入。')
            return result

        table_data = self._extract_pdf_tables(file_path)

        items = []
        if table_data:
            items = self._parse_items_from_table(table_data)

        text_items = self._parse_items_from_text(text)
        items = self._merge_items(items, text_items)

        result['items'] = items
        result['success'] = len(items) > 0

        if not items:
            result['errors'].append('未能从PDF中识别出持仓信息。')

        if client_id and items:
            self.save_to_repository(client_id, items, source='pdf',
                                    source_path=file_path)

        return result

    def _extract_pdf_text(self, file_path: str) -> str:
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                texts = []
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        texts.append(t)
                return '\n'.join(texts)
        except ImportError:
            pass
        except Exception:
            pass

        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            texts = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
            return '\n'.join(texts)
        except ImportError:
            pass
        except Exception:
            pass

        return ""

    def _extract_pdf_tables(self, file_path: str) -> list:
        try:
            import pdfplumber
            tables = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_tables = page.extract_tables()
                    for t in page_tables:
                        if t:
                            tables.extend(t)
            return tables
        except ImportError:
            return []
        except Exception:
            return []

    # ==================== 浏览器链接爬取 ====================

    def import_from_url(self, url: str, client_id: str = None,
                        credentials: dict = None, headless: bool = True) -> dict:
        """从浏览器链接爬取持仓信息"""
        result = {'success': False, 'items': [], 'errors': [], 'screenshot_path': ''}

        platform = self._detect_platform(url)

        if platform == 'eastmoney':
            return self._crawl_eastmoney(url, client_id, credentials, headless)
        elif platform == 'antfortune':
            return self._crawl_antfortune_hint()
        elif platform == 'xueqiu':
            return self._crawl_xueqiu_hint()

        return self._crawl_with_selenium(url, client_id, credentials, headless)

    def _detect_platform(self, url: str) -> str:
        url_lower = url.lower()
        if 'eastmoney.com' in url_lower or '天天基金' in url_lower:
            return 'eastmoney'
        if 'xueqiu.com' in url_lower or '雪球' in url_lower:
            return 'xueqiu'
        if 'antfortune.com' in url_lower or '支付宝' in url_lower:
            return 'antfortune'
        if 'howbuy.com' in url_lower:
            return 'howbuy'
        return 'unknown'

    def _crawl_eastmoney(self, url: str, client_id: str = None,
                         credentials: dict = None, headless: bool = True) -> dict:
        """爬取东方财富/天天基金持仓"""
        result = {'success': False, 'items': [], 'errors': [], 'screenshot_path': ''}

        # 优先用公开API
        codes = re.findall(r'(\d{6})', url)
        if codes and not credentials:
            return self._crawl_eastmoney_api(codes, client_id)

        # Selenium 方案
        try:
            from selenium import webdriver
            from selenium.webdriver.common.by import By
            from selenium.webdriver.chrome.options import Options
        except ImportError:
            return self._crawl_eastmoney_api(codes, client_id) if codes else result

        try:
            options = Options()
            if headless:
                options.add_argument('--headless=new')
            options.add_argument('--no-sandbox')
            options.add_argument('--disable-dev-shm-usage')
            options.add_argument('--window-size=1920,1080')

            driver = webdriver.Chrome(options=options)
            driver.get(url)

            from selenium.webdriver.support.ui import WebDriverWait
            from selenium.webdriver.support import expected_conditions as EC

            # 登录
            if credentials and credentials.get('username'):
                self._attempt_selenium_login(driver, credentials)

            # 截图
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            client_dir = self.clients_dir / (client_id or 'unknown') / 'imports'
            os.makedirs(str(client_dir), exist_ok=True)
            screenshot_path = str(client_dir / f'screenshot_{timestamp}.png')
            driver.save_screenshot(screenshot_path)
            result['screenshot_path'] = screenshot_path

            page_text = driver.find_element(By.TAG_NAME, 'body').text
            driver.quit()

            items = self._parse_items_from_text(page_text)
            result['items'] = items
            result['success'] = len(items) > 0

            if client_id and items:
                self.save_to_repository(client_id, items, source='url_crawl',
                                        source_path=url)

        except Exception as e:
            result['errors'].append(f'浏览器爬取失败: {str(e)}')

        return result

    def _crawl_eastmoney_api(self, codes: list, client_id: str = None) -> dict:
        """通过API抓取东方财富公开数据"""
        import sys
        sys.path.insert(0, str(SKILL_DIR / "pkg"))
        from crawl_utils import safe_request

        result = {'success': False, 'items': [], 'errors': [], 'screenshot_path': ''}
        items = []

        for code in codes:
            try:
                # 腾讯行情API (股票)
                ts_code = f"sh{code}" if code.startswith(('6', '5')) else f"sz{code}"
                url = f"https://qt.gtimg.cn/q={ts_code}"
                raw = safe_request(url, timeout=8)
                if raw:
                    text = raw.decode("gbk", errors="replace")
                    m = re.search(r'v_\w+="(.+?)"', text)
                    if m:
                        fields = m.group(1).split("~")
                        if len(fields) > 31:
                            items.append({
                                'code': code,
                                'name': fields[1],
                                'item_type': 'stock',
                                'price': self._sf(fields[3]),
                                'change_pct': self._sf(fields[31]),
                            })
                            continue

                # 天天基金API (基金)
                furl = f"https://fundgz.1234567.com.cn/js/{code}.js"
                raw = safe_request(furl, timeout=8)
                if raw:
                    text = raw.decode("utf-8", errors="replace")
                    m = re.search(r'\((.+)\)', text)
                    if m:
                        d = json.loads(m.group(1))
                        items.append({
                            'code': code,
                            'name': d.get('name', ''),
                            'item_type': 'fund',
                            'nav': self._sf(d.get('dwjz', 0)),
                            'change_pct': self._sf(d.get('gszzl', 0)),
                        })
            except Exception:
                pass

        result['items'] = items
        result['success'] = len(items) > 0
        return result

    def _crawl_antfortune_hint(self) -> dict:
        return {'success': False, 'items': [], 'errors': [
            '蚂蚁财富需要扫码登录，建议使用截图导入功能。请截屏持仓页面后使用 import_from_screenshot 导入。'
        ]}

    def _crawl_xueqiu_hint(self) -> dict:
        return {'success': False, 'items': [], 'errors': [
            '雪球需要Cookie认证登录，建议使用截图导入功能。'
        ]}

    def _crawl_with_selenium(self, url: str, client_id: str = None,
                             credentials: dict = None, headless: bool = True) -> dict:
        """通用Selenium爬取"""
        result = {'success': False, 'items': [], 'errors': [], 'screenshot_path': ''}

        try:
            from selenium import webdriver
            from selenium.webdriver.common.by import By
            from selenium.webdriver.chrome.options import Options
        except ImportError:
            result['errors'].append(
                '需要安装 Selenium: pip install selenium\n'
                '并安装 ChromeDriver 到系统 PATH。\n'
                '也可以将页面截图后使用截图导入功能。'
            )
            return result

        try:
            options = Options()
            if headless:
                options.add_argument('--headless=new')
            options.add_argument('--no-sandbox')
            options.add_argument('--disable-dev-shm-usage')
            options.add_argument('--window-size=1920,1080')

            driver = webdriver.Chrome(options=options)
            driver.get(url)

            if credentials and credentials.get('username'):
                self._attempt_selenium_login(driver, credentials)

            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            client_dir = self.clients_dir / (client_id or 'unknown') / 'imports'
            os.makedirs(str(client_dir), exist_ok=True)
            screenshot_path = str(client_dir / f'screenshot_{timestamp}.png')
            driver.save_screenshot(screenshot_path)
            result['screenshot_path'] = screenshot_path

            page_text = driver.find_element(By.TAG_NAME, 'body').text
            driver.quit()

            items = self._parse_items_from_text(page_text)
            result['items'] = items
            result['success'] = len(items) > 0

            if client_id and items:
                self.save_to_repository(client_id, items, source='url_crawl',
                                        source_path=url)

        except Exception as e:
            result['errors'].append(f'浏览器爬取失败: {str(e)}')

        return result

    def _attempt_selenium_login(self, driver, credentials: dict) -> bool:
        try:
            from selenium.webdriver.common.by import By

            username = credentials.get('username', '')
            password = credentials.get('password', '')

            inputs = driver.find_elements(By.TAG_NAME, 'input')
            text_inputs = [i for i in inputs if i.get_attribute('type') in ('text', 'email', 'tel', None)]
            password_inputs = [i for i in inputs if i.get_attribute('type') == 'password']

            if text_inputs and password_inputs:
                text_inputs[0].send_keys(username)
                password_inputs[0].send_keys(password)

                buttons = driver.find_elements(By.TAG_NAME, 'button')
                submit_btns = [b for b in buttons
                               if '登录' in (b.text or '') or 'login' in (b.get_attribute('class') or '').lower()]
                if submit_btns:
                    submit_btns[0].click()
                    return True
            return False
        except Exception:
            return False

    @staticmethod
    def _sf(v, default=0.0):
        try:
            return float(v) if v not in ("", "-", None, "N/A", "null") else default
        except (TypeError, ValueError):
            return default

    # ==================== 文本解析引擎 ====================

    def _parse_items_from_text(self, text: str) -> list:
        """从文本解析股票/基金持仓"""
        items = []
        if not text:
            return items

        lines = text.strip().split('\n')
        for line in lines:
            codes = self.CODE_PATTERN.findall(line)
            if not codes:
                continue
            for code in codes:
                item = self._extract_item_from_line(line, code)
                if item:
                    items.append(item)

        # 全文搜索未匹配的代码
        all_codes = self.CODE_PATTERN.findall(text)
        seen_codes = {i.get('code') for i in items}
        for code in all_codes:
            if code not in seen_codes:
                idx = text.find(code)
                context = text[max(0, idx - 100):min(len(text), idx + 100)]
                item = self._extract_item_from_line(context, code)
                if item:
                    items.append(item)

        return [self._normalize_item(i) for i in items]

    def _extract_item_from_line(self, line: str, code: str) -> Optional[dict]:
        """从一行文本提取单条持仓"""
        item = {'code': code}

        # 判断类型: 股票代码 (6开头=上海, 0/3开头=深圳, 5开头=上海ETF), 基金代码
        if code.startswith(('6', '0', '3', '5')):
            item['item_type'] = 'stock'
        else:
            item['item_type'] = 'fund'

        # 名称: 从代码后提取，贪婪匹配到下一个数字或关键字为止
        name_pats = [
            # 代码后面紧跟的名称: "600519 贵州茅台" -> "贵州茅台"
            rf'{code}\s*[\)）]?\s*([^\d\s,，;；\n\r]{{2,30}}?)\s*(?:[\d（(]|份|股|成本|买入|止损|止盈|目标|金额|$)',
            # 名称在代码前: "贵州茅台(600519)"
            rf'([^\d\s,，;；\n\r]{{2,30}})\s*[\(（]?\s*{code}',
            # 包含常见公司后缀的股票名
            r'([一-龥A-Za-z]{2,20}(?:控股|股份|集团|科技|医药|银行|保险|证券|汽车|电子|通信|食品|饮料|家电|茅台|时代|能源|地产|消费|成长|稳健))',
            # 包含常见基金后缀
            r'([一-龥A-Za-z]{2,20}(?:基金|混合|债券|指数|货币|ETF|LOF|联接|增强|灵活|量化))',
            # 最后兜底: 代码后2-8个中英文字符
            rf'{code}\s*([一-龥A-Za-z]{{2,16}})',
        ]
        for pat in name_pats:
            m = re.search(pat, line)
            if m:
                name = m.group(1).strip()
                name = re.sub(r'[\(\)（）\*]', '', name).strip()
                if len(name) >= 2:
                    item['name'] = name
                    break

        # 股数/份额
        shares_matches = [
            re.search(r'(\d+(?:,\d{3})*(?:\.\d+)?)\s*[股分份]', line),
            re.search(r'[股分份]\s*[:：]?\s*(\d+(?:,\d{3})*(?:\.\d+)?)', line),
            re.search(r'(?:持有|持仓)\s*(\d+(?:,\d{3})*(?:\.\d+)?)', line),
        ]
        for m in shares_matches:
            if m:
                shares_str = m.group(1).replace(',', '')
                item['shares'] = float(shares_str)
                break

        # 成本
        cost_matches = [
            re.search(r'(?:成本|买入价|均价|价格)[:：\s]*(\d+\.?\d{0,4})', line),
            re.search(r'(\d+\.\d{2,4})\s*(?:元|买入|成本)', line),
        ]
        for m in cost_matches:
            if m:
                item['cost'] = float(m.group(1))
                break

        # 金额
        amt_match = self.AMOUNT_PATTERN.search(line)
        if amt_match:
            item['amount'] = float(amt_match.group(1).replace(',', ''))

        # 日期
        date_match = self.DATE_PATTERN.search(line)
        if date_match:
            y, mo, d = date_match.group(1), date_match.group(2).zfill(2), (date_match.group(3) or '01').zfill(2)
            item['purchase_date'] = f'{y}-{mo}-{d}'

        # 止盈止损
        sl_match = self.STOP_LOSS_PATTERN.search(line)
        if sl_match:
            item['stop_loss'] = -abs(float(sl_match.group(1)))

        tp_match = self.TAKE_PROFIT_PATTERN.search(line)
        if tp_match:
            item['take_profit'] = float(tp_match.group(1))

        return item if len(item) > 2 else None

    def _parse_items_from_table(self, table_data: list) -> list:
        """从表格解析持仓"""
        if not table_data or len(table_data) < 2:
            return []

        header_row = None
        col_map = {}

        for i, row in enumerate(table_data):
            row_text = ' '.join(str(c) for c in row).lower()
            if any(kw in row_text for kw in ['代码', '名称', '份额', '成本', '净值', '股数', '金额', '日期', '类型']):
                header_row = i
                col_map = self._map_table_columns(row)
                break

        if not col_map:
            items = []
            for row in table_data:
                row_text = ' '.join(str(c) for c in row)
                codes = self.CODE_PATTERN.findall(row_text)
                for code in codes:
                    item = self._extract_item_from_line(row_text, code)
                    if item:
                        items.append(item)
            return items

        items = []
        for row in table_data[header_row + 1:]:
            if not row or all(not c or str(c).strip() == '' for c in row):
                continue

            item = {}
            for col_idx, field in col_map.items():
                if col_idx < len(row) and row[col_idx]:
                    val = str(row[col_idx]).strip()
                    if field == 'code':
                        match = self.CODE_PATTERN.search(val)
                        if match:
                            item['code'] = match.group(1)
                    elif field == 'name':
                        item['name'] = val
                    elif field in ('shares', 'cost', 'amount', 'price'):
                        try:
                            item[field] = float(val.replace(',', ''))
                        except ValueError:
                            pass
                    elif field == 'item_type':
                        item['item_type'] = 'stock' if '股' in val.lower() or 'stock' in val.lower() else 'fund'
                    elif field in ('stop_loss', 'take_profit'):
                        try:
                            item[field] = float(val.replace(',', '').replace('%', ''))
                        except ValueError:
                            pass
                    elif field == 'purchase_date':
                        item['purchase_date'] = val

            if item.get('code'):
                item.setdefault('item_type',
                               'stock' if item['code'].startswith(('6', '0', '3', '5')) else 'fund')
                items.append(item)

        return items

    def _map_table_columns(self, header_row: list) -> dict:
        col_map = {}
        for i, col in enumerate(header_row):
            col_lower = str(col).lower().strip()
            if any(kw in col_lower for kw in ['代码', 'code', '编号']):
                col_map[i] = 'code'
            elif any(kw in col_lower for kw in ['名称', 'name', '股票', '基金', '产品']):
                col_map[i] = 'name'
            elif any(kw in col_lower for kw in ['类型', 'type', '品种']):
                col_map[i] = 'item_type'
            elif any(kw in col_lower for kw in ['股数', '份额', '份数', 'shares', '数量', '持仓']):
                col_map[i] = 'shares'
            elif any(kw in col_lower for kw in ['成本', '买入价', 'cost', '均价', '价格']):
                col_map[i] = 'cost'
            elif any(kw in col_lower for kw in ['金额', '市值', 'amount', 'value', '总价']):
                col_map[i] = 'amount'
            elif any(kw in col_lower for kw in ['日期', 'date', '时间', '买入']):
                col_map[i] = 'purchase_date'
            elif any(kw in col_lower for kw in ['止损', 'stop_loss']):
                col_map[i] = 'stop_loss'
            elif any(kw in col_lower for kw in ['止盈', '目标', 'take_profit', 'target']):
                col_map[i] = 'take_profit'
        return col_map

    def _normalize_item(self, item: dict) -> dict:
        code = item.get('code', '')
        normalized = {
            'code': code,
            'item_type': item.get('item_type', 'stock' if code.startswith(('6', '0', '3', '5')) else 'fund'),
            'name': item.get('name', ''),
            'shares': item.get('shares'),
            'cost': item.get('cost'),
            'amount': item.get('amount'),
            'price': item.get('price') or item.get('nav'),
            'change_pct': item.get('change_pct'),
            'purchase_date': item.get('purchase_date'),
            'stop_loss': item.get('stop_loss'),
            'take_profit': item.get('take_profit'),
            'source': item.get('source', 'import'),
            'import_date': datetime.now().strftime('%Y-%m-%d %H:%M'),
        }
        return {k: v for k, v in normalized.items() if v is not None}

    def _merge_items(self, existing: list, new: list) -> list:
        seen = set()
        merged = []
        for item in existing + new:
            key = item.get('code', '')
            if key and key not in seen:
                seen.add(key)
                merged.append(item)
        return merged

    def validate_items(self, items: list) -> tuple:
        valid = []
        errors = []
        for i, item in enumerate(items):
            row_errs = []
            code = item.get('code', '')
            if not code or not re.match(r'^\d{6}$', code):
                row_errs.append(f'第{i+1}行: 代码无效 ({code})')
            if item.get('shares') is not None and item['shares'] <= 0:
                row_errs.append(f'第{i+1}行: 股数/份额必须大于0')
            if row_errs:
                errors.extend(row_errs)
            else:
                valid.append(item)
        return valid, errors

    # ==================== Excel/CSV 导出 ====================

    def export_to_excel(self, items: list, output_path: str = None,
                        client_id: str = None) -> str:
        """导出持仓到Excel"""
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
        except ImportError:
            raise ImportError('请安装 openpyxl: pip install openpyxl')

        if not output_path:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_dir = self.clients_dir / (client_id or 'export') / 'reports'
            os.makedirs(str(output_dir), exist_ok=True)
            output_path = str(output_dir / f'holdings_{timestamp}.xlsx')

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = '持仓明细'

        header_font = Font(name='微软雅黑', bold=True, size=11, color='FFFFFF')
        header_fill = PatternFill(start_color='2F5496', end_color='2F5496', fill_type='solid')
        header_align = Alignment(horizontal='center', vertical='center')
        cell_align = Alignment(horizontal='center', vertical='center')
        thin_border = Border(
            left=Side(style='thin'), right=Side(style='thin'),
            top=Side(style='thin'), bottom=Side(style='thin')
        )

        headers = ['序号', '类型', '代码', '名称', '持有数量', '成本价', '当前价',
                   '市值', '盈亏', '收益率%', '止盈', '止损', '来源']
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=h)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_align
            cell.border = thin_border

        red_font = Font(color='FF0000')
        green_font = Font(color='008000')

        total_market = 0
        total_cost_val = 0

        for i, item in enumerate(items):
            row = i + 2
            shares = item.get('shares', 0) or 0
            cost = item.get('cost', 0) or 0
            price = item.get('price') or item.get('nav') or cost

            market_value = shares * price
            total_cost = shares * cost
            profit = market_value - total_cost
            profit_pct = (profit / total_cost * 100) if total_cost > 0 else 0
            total_market += market_value
            total_cost_val += total_cost

            type_label = '股票' if item.get('item_type') == 'stock' else '基金'
            data = [
                i + 1, type_label,
                item.get('code', ''),
                item.get('name', ''),
                shares,
                round(cost, 4),
                round(price, 4),
                round(market_value, 2),
                round(profit, 2),
                round(profit_pct, 2),
                item.get('take_profit', ''),
                item.get('stop_loss', ''),
                item.get('source', '导入'),
            ]
            for col, val in enumerate(data, 1):
                cell = ws.cell(row=row, column=col, value=val)
                cell.alignment = cell_align
                cell.border = thin_border

            if profit > 0:
                ws.cell(row=row, column=9).font = red_font
                ws.cell(row=row, column=10).font = red_font
            elif profit < 0:
                ws.cell(row=row, column=9).font = green_font
                ws.cell(row=row, column=10).font = green_font

        # 汇总行
        sr = len(items) + 2
        total_profit = total_market - total_cost_val
        total_pct = (total_profit / total_cost_val * 100) if total_cost_val > 0 else 0
        summary_row = ['合计', '', '', '', '', '', '',
                       round(total_market, 2), round(total_profit, 2), round(total_pct, 2), '', '', '']
        for col, val in enumerate(summary_row, 1):
            cell = ws.cell(row=sr, column=col, value=val)
            cell.font = Font(bold=True) if col == 1 else cell.font
            cell.border = thin_border

        cols_width = [6, 7, 12, 18, 12, 10, 10, 14, 12, 10, 8, 8, 10]
        for col, w in enumerate(cols_width, 1):
            ws.column_dimensions[get_column_letter(col)].width = w
        ws.freeze_panes = 'A2'

        # 汇总sheet
        ws2 = wb.create_sheet('持仓汇总')
        ws2.cell(row=1, column=1, value='项目').font = Font(bold=True, name='微软雅黑')
        ws2.cell(row=1, column=2, value='数值').font = Font(bold=True, name='微软雅黑')
        summary_data = [
            ('持仓数量', f'{len(items)}条'),
            ('总投入成本', f'¥{total_cost_val:,.2f}'),
            ('当前总市值', f'¥{total_market:,.2f}'),
            ('总盈亏', f'¥{total_profit:,.2f}'),
            ('总收益率', f'{total_pct:+.2f}%'),
            ('导出时间', datetime.now().strftime('%Y-%m-%d %H:%M:%S')),
        ]
        for r, (label, val) in enumerate(summary_data, 2):
            ws2.cell(row=r, column=1, value=label).font = Font(bold=True, name='微软雅黑')
            ws2.cell(row=r, column=2, value=val).font = Font(name='微软雅黑')

        ws2.column_dimensions['A'].width = 18
        ws2.column_dimensions['B'].width = 25

        wb.save(output_path)
        return output_path

    def export_to_csv(self, items: list, output_path: str = None,
                      client_id: str = None) -> str:
        """导出持仓到CSV"""
        if not output_path:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_dir = self.clients_dir / (client_id or 'export') / 'reports'
            os.makedirs(str(output_dir), exist_ok=True)
            output_path = str(output_dir / f'holdings_{timestamp}.csv')

        fieldnames = [
            'item_type', 'code', 'name', 'shares', 'cost', 'price',
            'market_value', 'profit', 'profit_pct',
            'stop_loss', 'take_profit', 'purchase_date', 'source', 'import_date'
        ]

        with open(output_path, 'w', newline='', encoding='utf-8-sig') as f:
            writer = csv.DictWriter(f, fieldnames=fieldnames, extrasaction='ignore')
            writer.writeheader()

            for item in items:
                shares = item.get('shares', 0) or 0
                cost = item.get('cost', 0) or 0
                price = item.get('price') or item.get('nav') or cost
                market_value = shares * price
                total_cost = shares * cost
                profit = market_value - total_cost

                row = {
                    'item_type': item.get('item_type', 'stock'),
                    'code': item.get('code', ''),
                    'name': item.get('name', ''),
                    'shares': shares,
                    'cost': round(cost, 4),
                    'price': round(price, 4),
                    'market_value': round(market_value, 2),
                    'profit': round(profit, 2),
                    'profit_pct': round(profit / total_cost * 100, 2) if total_cost > 0 else 0,
                    'stop_loss': item.get('stop_loss', ''),
                    'take_profit': item.get('take_profit', ''),
                    'purchase_date': item.get('purchase_date', ''),
                    'source': item.get('source', 'import'),
                    'import_date': item.get('import_date', datetime.now().strftime('%Y-%m-%d')),
                }
                writer.writerow(row)

        return output_path

    # ==================== 客户持仓仓库 ====================

    def save_to_repository(self, client_id: str, items: list,
                           source: str = 'manual', source_path: str = None) -> str:
        """保存持仓到客户仓库并同步到跟踪器"""
        client_dir = self.clients_dir / client_id
        for sub in ['imports', 'reports', 'history']:
            (client_dir / sub).mkdir(parents=True, exist_ok=True)

        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        for item in items:
            item.setdefault('source', source)
            item.setdefault('import_date', datetime.now().strftime('%Y-%m-%d %H:%M'))
            item['import_batch'] = timestamp

        # 读取现有持仓并合并
        holdings_file = client_dir / 'holdings.json'
        existing = {}
        if holdings_file.exists():
            try:
                with open(holdings_file, 'r', encoding='utf-8') as f:
                    for h in json.load(f):
                        code = h.get('code', '')
                        if code:
                            existing[code] = h
            except Exception:
                pass

        for item in items:
            code = item.get('code', '')
            if code:
                existing[code] = item

        merged = list(existing.values())

        with open(holdings_file, 'w', encoding='utf-8') as f:
            json.dump(merged, f, ensure_ascii=False, indent=2)

        # 历史快照
        history_file = client_dir / 'history' / f'holdings_{timestamp}.json'
        with open(history_file, 'w', encoding='utf-8') as f:
            json.dump(merged, f, ensure_ascii=False, indent=2)

        # 导入记录
        import_log = client_dir / 'imports' / f'import_{timestamp}.json'
        with open(import_log, 'w', encoding='utf-8') as f:
            json.dump({
                'timestamp': timestamp, 'source': source,
                'source_path': source_path, 'count': len(items), 'items': items
            }, f, ensure_ascii=False, indent=2)

        # 同步到 SQLite tracker
        self._sync_to_sqlite_tracker(items)
        # 同步到 JSON tracker
        self._sync_to_json_tracker(items)

        self._update_repository_index(client_id, merged)
        return str(holdings_file)

    def _sync_to_sqlite_tracker(self, items: list):
        """同步到 SQLite PortfolioTracker"""
        tracker = self.sqlite_tracker
        if tracker is None:
            return
        for item in items:
            try:
                item_type = item.get('item_type', 'stock')
                code = item.get('code', '')
                name = item.get('name', '')
                cost = item.get('cost', 0) or 0
                shares = item.get('shares', 0) or 0
                tp = item.get('take_profit') or item.get('target_return', 0) or 0
                sl = item.get('stop_loss', -10)

                tracker.add_holding(
                    item_type=item_type, code=code, name=name,
                    cost=cost, shares=shares, target_return=tp
                )
                if sl != -10 or tp != 0:
                    tracker.update_settings(
                        item_type, code,
                        stop_loss=sl if sl != -10 else None,
                        stop_profit=tp if tp != 0 else None
                    )
            except Exception:
                pass

    def _sync_to_json_tracker(self, items: list):
        """同步到 JSON PortfolioTracker"""
        tracker = self.json_tracker
        if tracker is None:
            return
        for item in items:
            try:
                code = item.get('code', '')
                name = item.get('name', '')
                cost = item.get('cost')
                shares = item.get('shares')
                sl = item.get('stop_loss')
                tp = item.get('take_profit')
                tracker.add_stock(
                    code=code, name=name, cost=cost,
                    shares=shares, stop_loss=sl, take_profit=tp
                )
            except Exception:
                pass

    def load_client_items(self, client_id: str) -> list:
        holdings_file = self.clients_dir / client_id / 'holdings.json'
        if not holdings_file.exists():
            return []
        try:
            with open(holdings_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            return []

    def get_import_history(self, client_id: str) -> list:
        imports_dir = self.clients_dir / client_id / 'imports'
        if not imports_dir.exists():
            return []
        history = []
        for f in sorted(imports_dir.glob('import_*.json'), reverse=True):
            try:
                with open(f, 'r', encoding='utf-8') as fh:
                    history.append(json.load(fh))
            except Exception:
                pass
        return history

    def get_history_snapshots(self, client_id: str) -> list:
        history_dir = self.clients_dir / client_id / 'history'
        if not history_dir.exists():
            return []
        return [{'timestamp': f.stem.replace('holdings_', ''), 'path': str(f)}
                for f in sorted(history_dir.glob('holdings_*.json'), reverse=True)]

    def _update_repository_index(self, client_id: str, items: list):
        index_file = self.clients_dir / '_index.json'
        index = {}
        if index_file.exists():
            try:
                with open(index_file, 'r', encoding='utf-8') as f:
                    index = json.load(f)
            except Exception:
                pass

        total_value = 0
        for item in items:
            shares = item.get('shares', 0) or 0
            price = item.get('price') or item.get('nav') or item.get('cost', 0) or 0
            total_value += shares * price

        index[client_id] = {
            'name': index.get(client_id, {}).get('name', client_id),
            'items_count': len(items),
            'total_value': round(total_value, 2),
            'last_updated': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            'codes': [i.get('code') for i in items]
        }

        with open(index_file, 'w', encoding='utf-8') as f:
            json.dump(index, f, ensure_ascii=False, indent=2)

    def list_clients(self) -> dict:
        index_file = self.clients_dir / '_index.json'
        if index_file.exists():
            try:
                with open(index_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception:
                pass
        return {}

    def delete_client(self, client_id: str) -> bool:
        import shutil
        client_dir = self.clients_dir / client_id
        if client_dir.exists():
            shutil.rmtree(str(client_dir))
            self._update_repository_index(client_id, [])
            return True
        return False

    def print_summary(self, client_id: str = None):
        """打印持仓汇总 (集成量化跟踪)"""
        items = self.load_client_items(client_id) if client_id else []
        if not items:
            # 尝试从 SQLite tracker 获取
            if self.sqlite_tracker:
                summary = self.sqlite_tracker.get_portfolio_summary()
                items = summary.get('items', [])

        if not items:
            print("暂无持仓记录。")
            return

        print("\n" + "=" * 80)
        print(f"  持仓汇总" + (f" - {client_id}" if client_id else ""))
        print("=" * 80)
        print(f"  {'类型':<5} {'代码':<8} {'名称':<14} {'数量':>8} {'成本':>8} {'现价':>8} {'市值':>10} {'盈亏':>10} {'收益率':>8}")
        print("-" * 80)

        total_cost_val = 0
        total_market_val = 0
        for item in items:
            shares = item.get('shares', 0) or 0
            cost = item.get('cost', 0) or 0
            price = item.get('price') or item.get('current_price') or item.get('nav') or cost
            market = shares * price
            total_cost = shares * cost
            pnl = market - total_cost
            pnl_pct = (pnl / total_cost * 100) if total_cost > 0 else 0
            total_cost_val += total_cost
            total_market_val += market

            itype = item.get('item_type', 'stock')
            type_label = '股' if itype == 'stock' else '基'
            name = (item.get('name') or '')[:12]
            code = item.get('code', '')

            print(f"  {type_label:<5} {code:<8} {name:<14} {shares:>8.0f} {cost:>8.2f} {price:>8.2f} {market:>10.2f} {pnl:>+10.2f} {pnl_pct:>+7.2f}%")

        total_pnl = total_market_val - total_cost_val
        total_pct = (total_pnl / total_cost_val * 100) if total_cost_val > 0 else 0
        print("-" * 80)
        print(f"  合计: 成本 {total_cost_val:.2f}  市值 {total_market_val:.2f}  盈亏 {total_pnl:+.2f} ({total_pct:+.2f}%)")
        print("=" * 80)

        # 如有量化跟踪器，同时打印调仓建议
        if self.sqlite_tracker:
            print()
            self.sqlite_tracker.print_advice()


def main():
    """测试"""
    engine = StockImportEngine()
    print("=" * 60)
    print("  StockImportEngine v1.0 测试")
    print("=" * 60)

    # 测试文本解析
    test_text = """
    持仓明细
    600519 贵州茅台 1000股 成本1350.00 止损-10% 止盈20%
    000858 五粮液 5000股 成本95.00
    300750 宁德时代 2000股 成本450.00 目标25%
    110022 易方达消费行业 50000份 成本2.800
    """

    print("\n【文本解析测试】")
    items = engine._parse_items_from_text(test_text)
    for item in items:
        print(f"  [{item.get('item_type')}] {item.get('code')} {item.get('name', '?')} "
              f"shares={item.get('shares')} cost={item.get('cost')} "
              f"sl={item.get('stop_loss')} tp={item.get('take_profit')}")

    # 测试表格解析
    test_table = [
        ['代码', '名称', '类型', '持有数量', '成本价', '止损', '止盈'],
        ['600519', '贵州茅台', '股票', '1000', '1350', '-10', '20'],
        ['110022', '易方达消费', '基金', '50000', '2.8', '', '15'],
        ['000858', '五粮液', '股票', '5000', '95', '-8', ''],
    ]
    print("\n【表格解析测试】")
    table_items = engine._parse_items_from_table(test_table)
    for item in table_items:
        print(f"  [{item.get('item_type')}] {item.get('code')} {item.get('name', '?')} "
              f"shares={item.get('shares')} cost={item.get('cost')}")

    # 保存到仓库
    print("\n【仓库保存测试】")
    path = engine.save_to_repository('test_client', items, source='test')
    print(f"  保存到: {path}")
    loaded = engine.load_client_items('test_client')
    print(f"  加载: {len(loaded)} 条")

    # 导出CSV
    csv_path = engine.export_to_csv(items, client_id='test_client')
    print(f"\n【CSV导出】 {csv_path}")

    # 客户列表
    clients = engine.list_clients()
    print(f"\n【客户列表】 {list(clients.keys())}")

    print("\n" + "=" * 60)
    print("  测试完成")
    print("=" * 60)


if __name__ == '__main__':
    main()
