"""
PPT资产分析报告生成器 v1.0
生成可视化的基金资产分析报告（.pptx格式）
"""
import os
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PPTReportGenerator:
    """PPT资产分析报告生成器"""

    # 颜色配置
    COLORS = {
        'primary': (0, 102, 204),      # 蓝色
        'secondary': (51, 51, 51),     # 深灰
        'positive': (0, 128, 0),       # 绿色
        'negative': (204, 0, 0),       # 红色
        'neutral': (128, 128, 128),    # 灰色
        'background': (245, 245, 245),  # 浅灰背景
        'white': (255, 255, 255),
        'light_blue': (204, 229, 255),
    }

    def __init__(self, data_dir=None):
        self.data_dir = data_dir or Path(__file__).resolve().parent.parent.parent / "data"

    def generate_asset_analysis_ppt(self, user_id: str, holdings: list,
                                     analysis: dict = None, output_path: str = None) -> str:
        """
        生成PPT资产分析报告

        Args:
            user_id: 用户ID
            holdings: 持仓列表
            analysis: 分析数据（如果为None则实时获取）
            output_path: 输出路径（如果为None则自动生成）

        Returns:
            str: 生成的文件路径
        """
        if not HAS_PPTX:
            return "错误：python-pptx库未安装，请运行: pip install python-pptx"

        if not holdings:
            return "错误：没有持仓数据"

        # 创建输出目录
        if output_path is None:
            output_dir = self.data_dir / 'reports'
            output_dir.mkdir(parents=True, exist_ok=True)
            date_str = datetime.now().strftime('%Y-%m-%d_%H%M')
            output_path = str(output_dir / f'asset_report_{user_id}_{date_str}.pptx')

        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 生成各页
        self._add_title_slide(prs, user_id)
        self._add_overview_slide(prs, holdings, analysis)
        self._add_holdings_detail_slide(prs, holdings, analysis)
        self._add_allocation_slide(prs, holdings)
        self._add_summary_slide(prs, holdings, analysis)

        # 保存
        prs.save(output_path)
        return output_path

    def _add_title_slide(self, prs, user_id: str):
        """添加标题页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 背景色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(13.333), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(*self.COLORS['primary'])
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "基金资产分析报告"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RgbColor(*self.COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"用户: {user_id} | 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(*self.COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

    def _add_overview_slide(self, prs, holdings: list, analysis: dict):
        """添加总览页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 标题
        self._add_slide_title(slide, "资产总览", Inches(0.5), Inches(0.3), Pt(32))

        # 计算汇总数据
        total_value = sum(h.get('current_value', 0) for h in holdings)
        total_cost = sum(h.get('cost', 0) * h.get('shares', 0) for h in holdings)
        total_profit = total_value - total_cost
        profit_pct = (total_profit / total_cost * 100) if total_cost > 0 else 0

        # 汇总卡片
        cards = [
            ("总资产", f"{total_value:.2f}万元", self.COLORS['primary']),
            ("总成本", f"{total_cost:.2f}万元", self.COLORS['secondary']),
            ("总收益", f"{total_profit:+.2f}万元", self.COLORS['positive'] if total_profit >= 0 else self.COLORS['negative']),
            ("收益率", f"{profit_pct:+.2f}%", self.COLORS['positive'] if profit_pct >= 0 else self.COLORS['negative']),
        ]

        for i, (label, value, color) in enumerate(cards):
            left = Inches(0.5 + i * 3.1)
            self._add_card(slide, left, Inches(1.2), Inches(2.8), Inches(1.8), label, value, color)

        # 持仓数量
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.333), Inches(0.5))
        tf = info_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"持有基金数量: {len(holdings)} 只"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(*self.COLORS['secondary'])

        # 饼图区域标题
        chart_title = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(6), Inches(0.5))
        tf = chart_title.text_frame
        p = tf.paragraphs[0]
        p.text = "资产配置比例"
        p.font.size = Pt(20)
        p.font.bold = True

    def _add_holdings_detail_slide(self, prs, holdings: list, analysis: dict):
        """添加持仓明细页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._add_slide_title(slide, "持仓明细", Inches(0.5), Inches(0.3), Pt(32))

        # 表头
        headers = ["基金名称", "代码", "持有份额", "成本", "当前净值", "收益", "收益率"]
        col_widths = [Inches(2.5), Inches(1.2), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)]

        y_start = Inches(1.0)
        x_start = Inches(0.5)

        # 绘制表头背景
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x_start, y_start, sum(col_widths), Inches(0.5)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RgbColor(*self.COLORS['primary'])
        header_shape.line.fill.background()

        # 表头文字
        x = x_start
        for i, (header, width) in enumerate(zip(headers, col_widths)):
            tb = slide.shapes.add_textbox(x, y_start + Inches(0.1), width, Inches(0.35))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = header
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RgbColor(*self.COLORS['white'])
            p.alignment = PP_ALIGN.CENTER
            x += width

        # 数据行
        y = y_start + Inches(0.5)
        for row_idx, h in enumerate(holdings[:10]):  # 最多10行
            bg_color = RgbColor(*self.COLORS['light_blue']) if row_idx % 2 == 0 else RgbColor(*self.COLORS['white'])

            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x_start, y, sum(col_widths), Inches(0.45)
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = bg_color
            row_bg.line.fill.background()

            fund_name = h.get('fund_name', '')[:10]
            fund_code = h.get('fund_code', '')
            shares = h.get('shares', 0)
            cost = h.get('cost', 0)
            nav = h.get('current_nav', h.get('nav', 0))
            profit = h.get('profit', 0)
            profit_pct = h.get('profit_pct', 0)

            values = [
                fund_name, fund_code,
                f"{shares:.0f}",
                f"{cost:.3f}",
                f"{nav:.4f}" if nav else "-",
                f"{profit:+.2f}",
                f"{profit_pct:+.2f}%"
            ]

            x = x_start
            for val, width in zip(values, col_widths):
                tb = slide.shapes.add_textbox(x, y + Inches(0.1), width, Inches(0.35))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.text = str(val)
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.CENTER

                # 收益着色
                if val in [f"{profit:+.2f}", f"{profit_pct:+.2f}%"]:
                    if profit >= 0:
                        p.font.color.rgb = RgbColor(*self.COLORS['positive'])
                    else:
                        p.font.color.rgb = RgbColor(*self.COLORS['negative'])

                x += width

            y += Inches(0.45)

        # 说明
        note = slide.shapes.add_textbox(Inches(0.5), y + Inches(0.3), Inches(12), Inches(0.4))
        tf = note.text_frame
        p = tf.paragraphs[0]
        p.text = "注：数据为估算值，实际净值以基金公司披露为准"
        p.font.size = Pt(10)
        p.font.color.rgb = RgbColor(*self.COLORS['neutral'])

    def _add_allocation_slide(self, prs, holdings: list):
        """添加资产配置页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._add_slide_title(slide, "资产配置分析", Inches(0.5), Inches(0.3), Pt(32))

        # 简单的文字配置说明
        y = Inches(1.2)

        # 计算各基金占比
        total_value = sum(h.get('current_value', 0) for h in holdings)
        if total_value == 0:
            total_value = sum(h.get('cost', 0) * h.get('shares', 0) for h in holdings)

        for h in holdings[:8]:
            fund_name = h.get('fund_name', '')[:12]
            value = h.get('current_value', h.get('cost', 0) * h.get('shares', 0))
            pct = (value / total_value * 100) if total_value > 0 else 0

            # 基金名称和占比条
            name_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(3), Inches(0.4))
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{fund_name} ({h.get('fund_code', '')})"
            p.font.size = Pt(14)

            # 占比条
            bar_width = Inches(pct / 10)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(3.8), y + Inches(0.1), bar_width, Inches(0.3)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = RgbColor(*self.COLORS['primary'])
            bar.line.fill.background()

            # 百分比文字
            pct_box = slide.shapes.add_textbox(Inches(4 + bar_width.inches + 0.2), y, Inches(2), Inches(0.4))
            tf = pct_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{pct:.1f}%"
            p.font.size = Pt(14)
            p.font.bold = True

            y += Inches(0.7)

    def _add_summary_slide(self, prs, holdings: list, analysis: dict):
        """添加总结页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._add_slide_title(slide, "分析总结与建议", Inches(0.5), Inches(0.3), Pt(32))

        # 基本信息
        total_value = sum(h.get('current_value', 0) for h in holdings)
        total_cost = sum(h.get('cost', 0) * h.get('shares', 0) for h in holdings)
        total_profit = total_value - total_cost
        profit_pct = (total_profit / total_cost * 100) if total_cost > 0 else 0

        # 风险评估
        risk_level = "中等"
        if profit_pct < -15:
            risk_level = "较高（亏损较大）"
        elif profit_pct > 20:
            risk_level = "较低（收益较好）"

        suggestions = []

        if profit_pct < -10:
            suggestions.append("当前亏损较大，建议保持耐心，避免在底部割肉")
        elif profit_pct > 15:
            suggestions.append("收益表现良好，可考虑适当止盈部分仓位")
        else:
            suggestions.append("建议保持现有配置，关注市场结构性机会")

        if len(holdings) > 5:
            suggestions.append("持仓基金数量较多，建议精简聚焦核心标的")
        elif len(holdings) < 2:
            suggestions.append("持仓较为集中，建议适度分散降低风险")

        # 绘制内容
        y = Inches(1.2)

        # 基本情况卡片
        self._add_card(slide, Inches(0.5), y, Inches(3.8), Inches(1.5),
                      "风险等级", risk_level, self.COLORS['neutral'])

        # 建议内容
        advice_box = slide.shapes.add_textbox(Inches(0.5), y + Inches(1.8), Inches(12), Inches(0.5))
        tf = advice_box.text_frame
        p = tf.paragraphs[0]
        p.text = "投资建议："
        p.font.size = Pt(18)
        p.font.bold = True

        y = y + Inches(2.4)
        for i, suggestion in enumerate(suggestions, 1):
            sug_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(12), Inches(0.5))
            tf = sug_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{i}. {suggestion}"
            p.font.size = Pt(16)
            y += Inches(0.5)

        # 风险提示
        y += Inches(0.5)
        warning_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(12), Inches(0.8))
        tf = warning_box.text_frame
        p = tf.paragraphs[0]
        p.text = "⚠️ 风险提示：以上内容仅供参考，不构成投资建议。基金投资有风险，过往业绩不代表未来表现。"
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(*self.COLORS['neutral'])

        # 生成时间
        time_box = slide.shapes.add_textbox(Inches(9), Inches(6.8), Inches(4), Inches(0.4))
        tf = time_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"报告生成: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        p.font.size = Pt(10)
        p.font.color.rgb = RgbColor(*self.COLORS['neutral'])
        p.alignment = PP_ALIGN.RIGHT

    def _add_slide_title(self, slide, title: str, left, top, font_size):
        """添加幻灯片标题"""
        tb = slide.shapes.add_textbox(left, top, Inches(12), Inches(0.8))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = font_size
        p.font.bold = True
        p.font.color.rgb = RgbColor(*self.COLORS['primary'])

    def _add_card(self, slide, left, top, width, height, label, value, color):
        """添加卡片"""
        # 卡片背景
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(*self.COLORS['white'])
        shape.line.color.rgb = RgbColor(*color)
        shape.line.width = Pt(2)

        # 标签
        label_tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.4))
        tf = label_tb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(*self.COLORS['secondary'])
        p.alignment = PP_ALIGN.CENTER

        # 数值
        value_tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), Inches(0.8))
        tf = value_tb.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RgbColor(*color)
        p.alignment = PP_ALIGN.CENTER


def generate_ppt_report(user_id: str, holdings: list, analysis: dict = None, data_dir: str = None) -> str:
    """
    便捷函数：生成PPT资产分析报告

    Args:
        user_id: 用户ID
        holdings: 持仓列表，格式：[{'fund_code': '000001', 'fund_name': 'xxx', 'shares': 10000, 'cost': 1.5, 'current_nav': 1.6, 'profit': 1000, 'profit_pct': 6.67}]
        analysis: 分析数据（可选）
        data_dir: 数据目录（可选）

    Returns:
        str: 生成的文件路径或错误信息
    """
    generator = PPTReportGenerator(data_dir=data_dir)
    return generator.generate_asset_analysis_ppt(user_id, holdings, analysis)